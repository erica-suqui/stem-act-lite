"""
Generate a STEM-ACT Progress Report .pptx skeleton.
Run: python3 generate_progress_report.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ── Brand colors ─────────────────────────────────────────────────────────────
STEM_BLUE   = RGBColor(0x1A, 0x4A, 0x8A)   # deep blue – headings / title bg
STEM_TEAL   = RGBColor(0x00, 0x8B, 0x8B)   # accent
STEM_LIGHT  = RGBColor(0xF0, 0xF4, 0xFA)   # slide background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1C, 0x1C, 0x2E)
GRAY        = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def slide_background(slide, color=STEM_LIGHT):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, color)
    # push background behind everything else
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def accent_bar(slide, height=0.06, color=STEM_TEAL):
    add_rect(slide, 0, 0, 13.33, height, color)


def slide_header(slide, title, subtitle=None):
    """Dark header band with title + optional subtitle."""
    add_rect(slide, 0, 0, 13.33, 1.4, STEM_BLUE)
    accent_bar(slide, height=0.06, color=STEM_TEAL)
    add_text(slide, title, 0.4, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12.5, 0.45,
                 font_size=14, color=RGBColor(0xB0, 0xC8, 0xE8), italic=True)


def bullet_box(slide, items, left, top, width, height,
               font_size=16, bullet="•  "):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT


def placeholder_note(slide, left, top, width, height, label="[ TODO ]"):
    rect = add_rect(slide, left, top, width, height,
                    RGBColor(0xE8, 0xEF, 0xF8), line=True)
    rect.line.color.rgb = STEM_TEAL
    rect.line.width = Pt(1)
    add_text(slide, label, left + 0.1, top + (height / 2) - 0.15,
             width - 0.2, 0.4,
             font_size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title Slide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, STEM_BLUE)          # full blue bg
add_rect(slide, 0, 0, 13.33, 0.08, STEM_TEAL)          # top accent
add_rect(slide, 0, 7.42, 13.33, 0.08, STEM_TEAL)       # bottom accent

add_text(slide, "STEM-ACT Collaboration Hub",
         1, 1.5, 11.33, 1.2,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Progress Report",
         1, 2.8, 11.33, 0.8,
         font_size=26, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

add_text(slide, "CSC400  |  Spring 2026",
         1, 3.6, 11.33, 0.5,
         font_size=18, color=RGBColor(0x90, 0xA8, 0xC8), align=PP_ALIGN.CENTER)

add_text(slide, "Alexander Rodriguez  ·  Erica Suqui  ·  Joe Youn",
         1, 4.3, 11.33, 0.5,
         font_size=16, color=RGBColor(0x90, 0xA8, 0xC8), align=PP_ALIGN.CENTER)

add_text(slide, f"Client: Cheryl Tokarski – STEM-ACT Director",
         1, 5.1, 11.33, 0.4,
         font_size=14, color=RGBColor(0x70, 0x90, 0xB0), align=PP_ALIGN.CENTER)

add_text(slide, datetime.date.today().strftime("%B %Y"),
         1, 5.7, 11.33, 0.4,
         font_size=13, italic=True,
         color=RGBColor(0x60, 0x80, 0xA0), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Agenda
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Agenda")

agenda_items = [
    "1.  Project Overview & Goals",
    "2.  Work Completed",
    "3.  Live Demo  (optional)",
    "4.  Challenges & How We Addressed Them",
    "5.  Next Steps & Timeline",
    "6.  Q & A",
]
bullet_box(slide, agenda_items, 1.5, 1.7, 10, 5, font_size=20, bullet="")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Project Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Project Overview",
             subtitle="What is STEM-ACT Collaboration Hub?")

add_text(slide, "The Problem", 0.5, 1.6, 6, 0.4,
         font_size=16, bold=True, color=STEM_BLUE)
bullet_box(slide, [
    "Manual event submissions via Google Forms",
    "Copy-paste workflow burdens admins",
    "No structured approval process or partner accounts",
    "Inconsistent event data & no partner visibility",
], 0.5, 2.05, 5.8, 3, font_size=14)

add_text(slide, "Our Solution", 6.8, 1.6, 6, 0.4,
         font_size=16, bold=True, color=STEM_BLUE)
bullet_box(slide, [
    "Custom WordPress plugin replacing manual workflow",
    "Role-based accounts: Viewer, Partner, Admin, Super Admin",
    "Structured event submission & approval states",
    "Automated publishing once events are approved",
    "Partner dashboard for status tracking & notifications",
], 6.8, 2.05, 6, 3, font_size=14)

# divider line
add_rect(slide, 6.4, 1.55, 0.05, 3.5, STEM_TEAL)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Goals & Objectives
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Goals & Objectives")

goals = [
    ("Streamlined Submissions",
     "Partners register once and submit events through a structured form — no more repeated Google Forms."),
    ("Role-Based Access Control",
     "Four user roles (Viewer, Partner, Admin, Super Admin) with appropriate permissions at each level."),
    ("Approval Workflow",
     "Every event passes through Pending → Approved / Denied states with admin notes and notifications."),
    ("Automated Publishing",
     "Approved events are instantly published to the public-facing event catalog without manual copying."),
    ("Partner Dashboard",
     "Partners can track submission status, edit pending events, and receive email notifications."),
]

y = 1.65
for title, desc in goals:
    add_text(slide, f"✦  {title}", 0.6, y, 12, 0.32,
             font_size=15, bold=True, color=STEM_BLUE)
    add_text(slide, desc, 1.1, y + 0.32, 11.5, 0.38,
             font_size=13, color=GRAY)
    y += 0.85


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Work Completed
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Work Completed",
             subtitle="Key milestones delivered so far")

completed = [
    "SRS document drafted and reviewed with client (v1.0, Feb 2026)",
    "Database schema designed (users, partners, events, roles)",
    "Next.js frontend scaffold with dynamic partner & user cards",
    "Clickable filters on partner and user card components",
    "Docker-based local development environment configured",
    "Architecture diagram & component design completed",
    "[ ADD: any additional completed features ]",
]
bullet_box(slide, completed, 0.6, 1.55, 12, 5.5, font_size=15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Demo Slide (optional placeholder)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Live Demo  (Optional)")

add_text(slide, "Replace this slide with screenshots or a live walkthrough.",
         0.6, 1.6, 12, 0.5, font_size=16, italic=True, color=GRAY,
         align=PP_ALIGN.CENTER)

placeholder_note(slide, 0.8, 2.2, 11.5, 4.4,
                 label="[ Insert screenshots / screen recording here ]")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Challenges & How We Addressed Them
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Challenges & How We Addressed Them")

# Table-style two-column layout
add_rect(slide, 0.5, 1.55, 5.8, 0.45, STEM_BLUE)
add_text(slide, "Challenge", 0.6, 1.6, 5.6, 0.35,
         font_size=14, bold=True, color=WHITE)

add_rect(slide, 6.8, 1.55, 6, 0.45, STEM_BLUE)
add_text(slide, "How We Addressed It", 6.9, 1.6, 5.8, 0.35,
         font_size=14, bold=True, color=WHITE)

challenges = [
    ("[ Describe Challenge 1 ]",          "[ How it was resolved / mitigated ]"),
    ("[ Describe Challenge 2 ]",          "[ How it was resolved / mitigated ]"),
    ("[ Describe Challenge 3 ]",          "[ How it was resolved / mitigated ]"),
    ("[ Describe Challenge 4 ]",          "[ How it was resolved / mitigated ]"),
]

y = 2.1
for i, (ch, sol) in enumerate(challenges):
    bg = RGBColor(0xE8, 0xEF, 0xF8) if i % 2 == 0 else WHITE
    add_rect(slide, 0.5, y, 5.8, 0.7, bg)
    add_rect(slide, 6.8, y, 6, 0.7, bg)
    add_text(slide, ch,  0.65, y + 0.08, 5.5, 0.55, font_size=13)
    add_text(slide, sol, 6.95, y + 0.08, 5.7, 0.55, font_size=13)
    y += 0.72


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Next Steps & Timeline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
slide_header(slide, "Next Steps & Timeline")

add_text(slide, "Upcoming Milestones", 0.6, 1.55, 12, 0.4,
         font_size=16, bold=True, color=STEM_BLUE)

milestones = [
    ("Sprint X  (now → MM/DD)",
     "[ e.g., Complete partner registration form & backend API ]"),
    ("Sprint X+1  (MM/DD → MM/DD)",
     "[ e.g., Admin dashboard – approve/deny workflow ]"),
    ("Sprint X+2  (MM/DD → MM/DD)",
     "[ e.g., Email notifications & automated publishing ]"),
    ("Sprint X+3  (MM/DD → MM/DD)",
     "[ e.g., Testing, bug fixes, user acceptance testing with client ]"),
    ("Final Delivery  (MM/DD)",
     "[ e.g., Production deployment & handoff to STEM-ACT ]"),
]

y = 2.1
for sprint, desc in milestones:
    add_rect(slide, 0.5, y, 0.08, 0.55, STEM_TEAL)   # left accent
    add_text(slide, sprint, 0.75, y,       11, 0.28, font_size=14, bold=True, color=STEM_BLUE)
    add_text(slide, desc,   0.75, y + 0.28, 11, 0.28, font_size=13, color=GRAY)
    y += 0.75


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Q & A / Thank You
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, STEM_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, STEM_TEAL)
add_rect(slide, 0, 7.42, 13.33, 0.08, STEM_TEAL)

add_text(slide, "Questions & Discussion",
         1, 2.0, 11.33, 1.2,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Thank you for your time!",
         1, 3.3, 11.33, 0.6,
         font_size=22, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

add_text(slide, "Alexander Rodriguez  ·  Erica Suqui  ·  Joe Youn",
         1, 4.2, 11.33, 0.4,
         font_size=15, color=RGBColor(0x90, 0xA8, 0xC8), align=PP_ALIGN.CENTER)

add_text(slide, "Client: Cheryl Tokarski – STEM-ACT Director",
         1, 4.75, 11.33, 0.4,
         font_size=13, italic=True,
         color=RGBColor(0x70, 0x90, 0xB0), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/j/Documents/stem-act/Documentation/STEM-ACT_Progress_Report.pptx"
prs.save(out)
print(f"Saved → {out}")
